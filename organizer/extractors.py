from __future__ import annotations

import io
import logging
import os
import sys
import zipfile
from dataclasses import dataclass, field
from pathlib import Path

from PIL import Image

LOGGER = logging.getLogger(__name__)
SUPPORTED_EXTENSIONS = {".pdf", ".doc", ".docx", ".ppt", ".pptx"}
MAX_OCR_PAGES = 30
MAX_OCR_IMAGES = 30
MAX_SOURCE_FILE_BYTES = 512 * 1024 * 1024
MAX_PDF_PAGES = 500
MAX_OFFICE_UNCOMPRESSED_BYTES = 512 * 1024 * 1024
MAX_OFFICE_MEMBER_BYTES = 128 * 1024 * 1024
MAX_EMBEDDED_IMAGE_BYTES = 64 * 1024 * 1024
MAX_OCR_IMAGE_PIXELS = 40_000_000
OCR_TIMEOUT_SECONDS = 90
MSO_AUTOMATION_SECURITY_FORCE_DISABLE = 3


class ExtractionLimitError(RuntimeError):
    """文件超过本地安全处理预算，可由调用方记录并跳过。"""


@dataclass
class ExtractionResult:
    text: str = ""
    warnings: list[str] = field(default_factory=list)
    ocr_used: bool = False
    blocked_reason: str = ""


def configure_ocr_engine() -> None:
    """让 Windows 打包版优先使用随应用分发的 Tesseract。"""
    try:
        import pytesseract
    except ImportError:
        return
    roots = [Path(getattr(sys, "_MEIPASS", "")), Path(sys.executable).parent, Path(__file__).resolve().parents[1]]
    for root in roots:
        binary = root / "tesseract" / "tesseract.exe"
        if binary.is_file():
            pytesseract.pytesseract.tesseract_cmd = str(binary)
            os.environ.setdefault("TESSDATA_PREFIX", str(binary.parent / "tessdata"))
            return


def extract_document(path: str | Path) -> ExtractionResult:
    document = Path(path)
    try:
        if document.stat().st_size > MAX_SOURCE_FILE_BYTES:
            raise ExtractionLimitError(f"文件超过 {MAX_SOURCE_FILE_BYTES // (1024 * 1024)} MB 的安全处理上限。")
    except OSError as error:
        raise ExtractionLimitError(f"无法读取文件大小：{error}") from error
    suffix = document.suffix.lower()
    if suffix == ".pdf":
        return _extract_pdf(document)
    if suffix == ".docx":
        return _extract_docx(document)
    if suffix == ".doc":
        return _extract_legacy_doc(document)
    if suffix == ".pptx":
        return _extract_pptx(document)
    if suffix == ".ppt":
        return _extract_legacy_ppt(document)
    return ExtractionResult(warnings=[f"不支持的格式：{suffix}"])


def _extract_pdf(path: Path) -> ExtractionResult:
    try:
        import pypdf
        import pypdfium2 as pdfium
    except ImportError:
        return ExtractionResult(warnings=["PDF 组件尚未安装。"])
    try:
        reader = pypdf.PdfReader(path)
        if len(reader.pages) > MAX_PDF_PAGES:
            raise ExtractionLimitError(f"PDF 超过 {MAX_PDF_PAGES} 页，为避免长时间处理已停止读取。")
        chunks = [page.extract_text() or "" for page in reader.pages]
        text = "\n".join(chunks).strip()
        needs_ocr = len(text) < max(200, len(reader.pages) * 80)
        if not needs_ocr:
            return ExtractionResult(text=text)

        pdf = pdfium.PdfDocument(path)
        ocr_text: list[str] = []
        skipped_for_limit = 0
        for index, page in enumerate(pdf):
            if index >= MAX_OCR_PAGES:
                break
            width, height = page.get_size()
            if width * 2 * height * 2 > MAX_OCR_IMAGE_PIXELS:
                skipped_for_limit += 1
                continue
            bitmap = page.render(scale=2)
            found = _ocr_image(bitmap.to_pil())
            if found:
                ocr_text.append(found)
        if ocr_text:
            return ExtractionResult(text=(text + "\n" + "\n".join(ocr_text)).strip(), ocr_used=True)
        if skipped_for_limit and len(text) < 100:
            reason = f"已跳过 {skipped_for_limit} 页超过安全处理上限的 PDF 页面。"
            return ExtractionResult(text=text, warnings=[reason], blocked_reason=reason)
        return ExtractionResult(text=text, warnings=["此 PDF 文字很少，但本机 OCR 没有可用结果。"])
    except ExtractionLimitError:
        raise
    except Exception as error:
        LOGGER.exception("Could not extract PDF: %s", path)
        reason = f"无法读取 PDF：{error}"
        return ExtractionResult(warnings=[reason], blocked_reason=reason)


def _extract_docx(path: Path) -> ExtractionResult:
    try:
        from docx import Document
        _validate_office_archive(path)
        document = Document(path)
        parts = [paragraph.text for paragraph in document.paragraphs]
        for table in document.tables:
            for row in table.rows:
                parts.extend(cell.text for cell in row.cells)
        return _append_embedded_image_ocr(path, "word/media/", "\n".join(parts).strip())
    except ExtractionLimitError:
        raise
    except Exception as error:
        LOGGER.exception("Could not extract DOCX: %s", path)
        reason = f"无法读取 Word 文件：{error}"
        return ExtractionResult(warnings=[reason], blocked_reason=reason)


def _extract_legacy_doc(path: Path) -> ExtractionResult:
    """通过已安装的 Microsoft Word 读取旧 .doc，不要求用户手动转换。"""
    word = document = previous_security = None
    try:
        client = _office_client()
        word = client.DispatchEx("Word.Application")
        word.Visible = False
        word.DisplayAlerts = 0
        previous_security = _disable_office_macros(word)
        document = word.Documents.Open(str(path.resolve()), False, True, False)
        return ExtractionResult(text=(document.Content.Text or "").strip())
    except Exception as error:
        return ExtractionResult(warnings=[f"无法读取旧 Word .doc：{_office_help(error)}"])
    finally:
        if document is not None:
            try:
                document.Close(False)
            except Exception:
                pass
        if word is not None:
            try:
                if previous_security is not None:
                    word.AutomationSecurity = previous_security
                word.Quit()
            except Exception:
                pass


def _extract_pptx(path: Path) -> ExtractionResult:
    try:
        from pptx import Presentation
        _validate_office_archive(path)
        presentation = Presentation(path)
        parts: list[str] = []
        for slide in presentation.slides:
            for shape in slide.shapes:
                if getattr(shape, "has_text_frame", False):
                    parts.append(shape.text)
                if getattr(shape, "has_table", False):
                    for row in shape.table.rows:
                        parts.extend(cell.text for cell in row.cells)
        return _append_embedded_image_ocr(path, "ppt/media/", "\n".join(parts).strip())
    except ExtractionLimitError:
        raise
    except Exception as error:
        LOGGER.exception("Could not extract PPTX: %s", path)
        reason = f"无法读取 PowerPoint 文件：{error}"
        return ExtractionResult(warnings=[reason], blocked_reason=reason)


def _extract_legacy_ppt(path: Path) -> ExtractionResult:
    """通过已安装的 Microsoft PowerPoint 读取旧 .ppt 的文字内容。"""
    powerpoint = presentation = previous_security = None
    try:
        client = _office_client()
        powerpoint = client.DispatchEx("PowerPoint.Application")
        powerpoint.DisplayAlerts = 0
        previous_security = _disable_office_macros(powerpoint)
        presentation = powerpoint.Presentations.Open(str(path.resolve()), True, False, False)
        parts: list[str] = []
        for slide in presentation.Slides:
            for shape in slide.Shapes:
                if shape.HasTextFrame and shape.TextFrame.HasText:
                    parts.append(shape.TextFrame.TextRange.Text)
                if shape.HasTable:
                    for row in range(1, shape.Table.Rows.Count + 1):
                        for column in range(1, shape.Table.Columns.Count + 1):
                            parts.append(shape.Table.Cell(row, column).Shape.TextFrame.TextRange.Text)
        return ExtractionResult(text="\n".join(parts).strip())
    except Exception as error:
        return ExtractionResult(warnings=[f"无法读取旧 PowerPoint .ppt：{_office_help(error)}"])
    finally:
        if presentation is not None:
            try:
                presentation.Close()
            except Exception:
                pass
        if powerpoint is not None:
            try:
                if previous_security is not None:
                    powerpoint.AutomationSecurity = previous_security
                powerpoint.Quit()
            except Exception:
                pass


def _office_client():
    if os.name != "nt":
        raise RuntimeError("旧 Office 格式只能在 Windows 上读取")
    try:
        import win32com.client
        return win32com.client
    except ImportError as error:
        raise RuntimeError("缺少 Windows Office 读取组件") from error


def _disable_office_macros(application):
    """自动化打开旧 Office 文件前必须强制关闭宏。"""
    try:
        previous_security = application.AutomationSecurity
        application.AutomationSecurity = MSO_AUTOMATION_SECURITY_FORCE_DISABLE
        return previous_security
    except Exception as error:
        raise RuntimeError("无法为 Office 自动化显式禁用宏，已停止读取该文件。") from error


def _office_help(error: Exception) -> str:
    if os.name == "nt":
        return f"请确认本机已安装 Microsoft Word/PowerPoint，然后重试。（{error}）"
    return str(error)


def _validate_office_archive(path: Path) -> None:
    """在解压 DOCX/PPTX 前限制压缩包展开规模，避免异常文件耗尽资源。"""
    try:
        with zipfile.ZipFile(path) as archive:
            members = archive.infolist()
    except (OSError, zipfile.BadZipFile) as error:
        raise ExtractionLimitError(f"Office 文件压缩包无效：{error}") from error
    total_size = sum(member.file_size for member in members)
    if total_size > MAX_OFFICE_UNCOMPRESSED_BYTES:
        raise ExtractionLimitError(f"Office 文件解压后超过 {MAX_OFFICE_UNCOMPRESSED_BYTES // (1024 * 1024)} MB 的安全处理上限。")
    oversized_member = next((member for member in members if member.file_size > MAX_OFFICE_MEMBER_BYTES), None)
    if oversized_member:
        raise ExtractionLimitError(f"Office 文件含有异常大的内部文件：{oversized_member.filename}。")


def _append_embedded_image_ocr(path: Path, prefix: str, text: str) -> ExtractionResult:
    if len(text) >= 300:
        return ExtractionResult(text=text)
    ocr_text: list[str] = []
    skipped_for_limit = 0
    try:
        with zipfile.ZipFile(path) as archive:
            images = [item for item in archive.infolist() if item.filename.startswith(prefix)]
            for image_info in images[:MAX_OCR_IMAGES]:
                try:
                    if image_info.file_size > MAX_EMBEDDED_IMAGE_BYTES:
                        skipped_for_limit += 1
                        continue
                    image = Image.open(io.BytesIO(archive.read(image_info)))
                    if image.width * image.height > MAX_OCR_IMAGE_PIXELS:
                        skipped_for_limit += 1
                        image.close()
                        continue
                    found = _ocr_image(image)
                    image.close()
                    if found:
                        ocr_text.append(found)
                except Exception:
                    LOGGER.debug("Skipping embedded image %s in %s", image_info.filename, path)
    except Exception:
        LOGGER.debug("Cannot inspect embedded images in %s", path, exc_info=True)
    if ocr_text:
        return ExtractionResult(text=(text + "\n" + "\n".join(ocr_text)).strip(), ocr_used=True)
    warning = "文字内容很少，已尝试读取内嵌图片，但本机 OCR 没有可用结果。"
    warnings = [warning] if len(text) < 100 else []
    if skipped_for_limit:
        limit_warning = f"已跳过 {skipped_for_limit} 张超过安全处理上限的内嵌图片。"
        warnings.append(limit_warning)
        if len(text) < 100:
            return ExtractionResult(text=text, warnings=warnings, blocked_reason=limit_warning)
    return ExtractionResult(text=text, warnings=warnings)


def _ocr_image(image: Image.Image) -> str:
    try:
        import pytesseract
        return pytesseract.image_to_string(image.convert("RGB"), lang="chi_sim+eng", timeout=OCR_TIMEOUT_SECONDS).strip()
    except Exception:
        LOGGER.debug("OCR unavailable", exc_info=True)
        return ""
